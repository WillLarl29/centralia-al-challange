from pathlib import Path

from pptx import Presentation

from .base import BaseLoader, RawDocument


class PptxLoader(BaseLoader):
    def load(self, path: Path) -> list[RawDocument]:
        presentation = Presentation(str(path))
        documents: list[RawDocument] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            texts.append(" | ".join(cells))

            content = "\n".join(texts).strip()
            if content:
                documents.append(
                    RawDocument(
                        content=content,
                        metadata={"source": path.name, "slide": slide_number, "type": "pptx"},
                    )
                )
        return documents
